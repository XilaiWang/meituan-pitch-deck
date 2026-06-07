#!/usr/bin/env python3
# 把渲染好的 16 页（assets/render/slide-NN.png, 2560x1440）装配为：
#  - 明信片-路演.pptx  整页满图 16:9 + 每页演讲者备注（WPS/PPT/Keynote 可开）
#  - 明信片-路演.pdf   投影兜底（Pillow 多页）
#  - 路演讲稿.md       逐页口播讲稿
import os, glob
from PIL import Image
from pptx import Presentation
from pptx.util import Emu

BASE = os.path.dirname(os.path.abspath(__file__))
REN = os.path.join(BASE, 'assets', 'render')
shots = sorted(glob.glob(os.path.join(REN, 'slide-*.png')))
assert len(shots) == 19, f'expected 19 slides, got {len(shots)}'

TITLES = [
 "封面 · 明信片", "问题 · 三难", "市场 · 3.2 亿", "答案 · 感知→托管→履约闭环",
 "入口 · 寄生在小团里", "信任地基 · 评估师建档", "看见 · 两张明信片", "洞察 · 还能吃 4 天",
 "调度 · 诚实度证明", "闭环 · 你帮了她", "回扣 · 双端联动", "双向 · 适老化与反向通道",
 "技术架构 · 1 Agent + 3 Skill", "商业模式 · 订阅分层", "成本对比 · 明信片 vs 保姆", "谁需要明信片 · 目标人群", "财务测算 · 三年模型",
 "伦理 · 三条红线 + 证据链", "收尾 · 一个画面",
]
NOTES = [
 "开场立住：明信片 · 美团 AI 居家养老管家。强调主角是现场真机 live demo，这份 deck 只作导引与投影兜底。",
 "三难层层递进：看不见、帮不上、像监控。点出第三难「监控」正是要打倒的敌人，也是差异化支点。监控这个词，全场只在这一页出现。",
 "市场只用核查过的数：3.23 亿、占 23.0%（国家统计局 2025 年末，已联网核实）。近 6 成空巢、慢病率约 75%，银发经济 7 万亿迈向 30 万亿。机会绑在美团履约能力上。",
 "一句立论：不是监控仪表盘，是会替你开口的明信片。情感入口→一键履约。第一次埋「为什么非美团不可」的四张网。",
 "入口：寄生在小团「为你推荐」的橙色卡，不必下新 App。副标题一句话立住商业模式：评估师上门 + Agent 长期跟。",
 "信任地基：人先上门建档、本人签字，AI 之后才学习。这是和算法监控的根本分界。可在此带出团队 / 照护师的可信度。",
 "看见：两类明信片物理隔离。情感卡只建议打电话、没有下单按钮；只有已有处方的药才出现「补一份」。绝不把喝粥和买药绑因果。",
 "洞察：还剩 4 天是按处方 + 服药打卡动态算出的确定结果，不是推荐。同样的算法下高血脂药 18 天就不打扰。只算余量、不诊断。",
 "调度·诚实：指着 trace 四态图例给评委看。药店是降级样例，real_api 尚未接入。我们选择诚实标注而非假装真实，这本身是工程成熟度。",
 "闭环：确认→模拟下单。演示标签和屏幕一样大方，不藏角注，不假装已对接库存。该有的都有，唯独不会真扣钱。",
 "情感高潮：此处真机停一拍。你这边点一份，2 秒后妈妈手机弹出绿色「小明关心你」。关怀第一次被她看见。",
 "双向 / 适老化：不是单向监控。妈妈也能一句话反向提需求，不会打字就按话筒。时间轴适老化、大字大按钮。",
 "技术架构：刻意放在情感高潮之后做可信度兜底。一个 OpenClaw Agent 串三个 Skill，AI 职责严格框在翻译 / 算余量 / 调度三件事。",
 "商业：订阅分层月费——基础 99 / 关怀 299（本次演示档，含上门）/ 生活 699 / 照护 1299 起；第二层走美团履约抽佣。为什么是美团：份额超七成、25 万药店。",
 "成本对比：一线城市住家保姆 6k–10k/月 vs 明信片关怀版 ¥299/月，省 80–95%。即使顶配照护版也仅为保姆费用的 15–25%。明信片不是替代保姆——是以 1/10 成本提供 24h AI 照护 + 美团全链路履约。",
 "谁需要明信片：空巢、能自理、但需要被看见的老人。请保姆太贵太早不划算，什么都不做又放不下。明信片填补「自理」和「全天陪护」之间的空白——¥299 先解决 80% 的焦虑。",
 "财务：商赛模型——3→10→30 城，付费 2 万→12 万→50 万，ARPU 420→600 元/月，第 3 年约 36 亿订阅；建档转化>25%、月留存>75%。强调是假设、可按城市供给成本调。",
 "伦理：三条红线都是屏幕原文不是承诺，随时可关、只补处方、模拟标注；服务过程靠 GPS 签到/节点打卡/图像查重/人工抽检的证据链。建议真机演示一次隐私 toggle。",
 "收尾：落在一个画面，妈妈在喝粥，你帮她补了一盒她日常在吃的降压药。克制、有温度、能落地。谢谢。",
]

# ── PPTX ──
prs = Presentation()
prs.slide_width = Emu(12192000)   # 13.333 in
prs.slide_height = Emu(6858000)   # 7.5 in
blank = prs.slide_layouts[6]
for i, p in enumerate(shots):
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(p, 0, 0, width=prs.slide_width, height=prs.slide_height)
    notes = s.notes_slide.notes_text_frame
    notes.text = f"S{i+1} {TITLES[i]}\n\n{NOTES[i]}"
pptx_path = os.path.join(BASE, '明信片-路演.pptx')
prs.save(pptx_path)

# ── PDF ──
imgs = [Image.open(p).convert('RGB') for p in shots]
pdf_path = os.path.join(BASE, '明信片-路演.pdf')
imgs[0].save(pdf_path, save_all=True, append_images=imgs[1:], resolution=144.0)

# ── 讲稿 ──
md = ["# 明信片 · 路演讲稿\n",
      "> 美团 2026 AI 黑客松。现场主角是真机 live demo，这份 deck 是导引与投影兜底。\n",
      "> 节奏建议：全场 5 到 7 分钟，情感高潮（S11）真机停一拍，技术与商业各 30 秒带过。\n"]
for i in range(19):
    md.append(f"\n## S{i+1} · {TITLES[i]}\n")
    md.append(f"{NOTES[i]}\n")
md.append("\n---\n\n## 评委 Q&A 速答\n")
qa = [
 ("这是真能跑还是 PPT 假数据？", "打开真机，指 trace 四态图例：药店是 fixture_fallback 降级样例、real_api 尚未接入、下单为 simulated；我们选择诚实标注 + 可视化 trace，而不是假装真实。"),
 ("药品相关，会不会变相荐药、有资质风险？", "AI 只补「已有处方」，数据来自持证照护师上门清点 + 本人签字（assessor_record）；不诊断、不推荐新药、不绑健康因果；余量是确定性算式得出；模拟下单不触碰真实药品交易。"),
 ("和老人监控 / 定位手环有何本质不同？", "三点：①情感入口（明信片）而非数据仪表盘；②妈妈主动授权、随时可关、子女看不到聊天和实时位置；③一键转美团真实履约（药 / 餐 / 家政）。"),
 ("订阅分层（99/299/699/1299 元/月）单位经济怎么算？", "见财务页：3→10→30 城、付费 2 万→12 万→50 万、ARPU 420→600 元/月，第 3 年约 36 亿订阅；建档一次性获客、季度回访摊薄、每次 care 走美团履约抽佣。是商赛模型假设，可按城市供给成本调，不编造。"),
 ("为什么是美团黑客松项目，而不是独立 App？", "唯一同时握有本地履约网络 + 即时配送 + 到家服务 + 小团 AI 入口的玩家；医药即时零售份额超七成、覆盖 25 万家药店，情感入口才能一键转真实履约。"),
]
for q,a in qa:
    md.append(f"\n**Q：{q}**\n\nA：{a}\n")
with open(os.path.join(BASE, '路演讲稿.md'), 'w', encoding='utf-8') as f:
    f.write("".join(md))

def mb(p): return f"{os.path.getsize(p)/1024/1024:.2f} MB"
print("PPTX :", pptx_path, mb(pptx_path))
print("PDF  :", pdf_path, mb(pdf_path))
print("讲稿 :", os.path.join(BASE, '路演讲稿.md'))
